# -*- coding: utf-8 -*-
from pptx import Presentation

SRC = "模板_副本.pptx"
OUT = "中期报告_v1文字版.pptx"
p = Presentation(SRC)

def find(shapes, name):
    for sh in shapes:
        if sh.shape_type == 6:
            r = find(sh.shapes, name)
            if r is not None: return r
        elif sh.name == name:
            return sh
    return None

def find_all(shapes, name, acc):
    for sh in shapes:
        if sh.shape_type == 6:
            find_all(sh.shapes, name, acc)
        elif sh.name == name:
            acc.append(sh)
    return acc

def fill(shape, text):
    """替换文字，保留首个 run 的字体/字号/颜色"""
    if shape is None: return False
    tf = shape.text_frame
    p0 = tf.paragraphs[0]
    if p0.runs:
        p0.runs[0].text = text
        for r in p0.runs[1:]:
            r._r.getparent().remove(r._r)
    else:
        p0.text = text
    for extra in tf.paragraphs[1:]:
        extra._p.getparent().remove(extra._p)
    return True

def do(slide_idx, mapping):
    s = p.slides[slide_idx]
    for name, text in mapping.items():
        sh = find(s.shapes, name)
        ok = fill(sh, text)
        if not ok:
            print(f"  ⚠ slide{slide_idx+1} 未找到 {name}")

# ---- S1 封面 ----
do(0, {
    "矩形 21": "面向云数据库高敏数据暴露路径侦测的证据约束智能体方法研究",
    "文本框 16": "2026 年 7 月    July 2026",
    "文本框 22": "王凌云　导师：蔡忠闽",
})
# ---- S2 目录 ----
do(1, {
    "矩形 44": "研究背景与研究目标",
    "矩形 35": "总体方案与技术路线",
    "矩形 28": "已完成工作",
    "矩形 40": "下一步工作计划",
})
# ---- S3 研究背景与问题(章节) ----
do(2, {
    "矩形 53": "研究背景与问题",
    "矩形 1": "高敏数据核心载体",
    "矩形 6": "暴露路径 Toxic Combination",
    "矩形 7": "判定需可审计",
})
# ---- S4 研究目标与总体方案 ----
do(3, {
    "矩形 22": "研究目标与总体方案",
    "矩形 17": "① CloudDB-PathBench 评测基准",
    "文本框 16": "面向暴露路径的受控、可复现评测环境",
    "矩形 19": "② EIC-Agent 证据约束智能体",
    "文本框 18": "表达—判定分离，判定可审计、抗幻觉",
    "矩形 21": "③ CloudDB-PathDetect 原型系统",
    "文本框 20": "端到端贯通与可视化验证",
})
# ---- S5 技术路线总览 ----
do(4, {
    "矩形 22": "技术路线总览",
    "矩形 24": "统一风险图 CDB-RG",
    "矩形 57": "身份·网络·数据·行为多源信号收编为同构图，作为智能体可查询环境",
    "矩形 44": "EIC-Agent 工具增强智能体",
    "矩形 58": "工具化取证 + 表达—判定分离，判定交由确定性算子",
    "矩形 51": "提前终止 + GV-FA 延伸",
    "矩形 63": "Gate 一票否决剪枝；图验证反馈对齐（技术路线已定）",
})
# ---- S6 关键实测:提前终止 ----
do(5, {
    "矩形 70": "关键实测：提前终止效率",
    "矩形 54": "≈25%",
    "矩形 29": "路径提前终止比例",
    "矩形 31": "68 条候选中 17 条被硬门一票否决剪枝",
    "矩形 58": "10×",
    "矩形 59": "被终止路径开销下降",
    "矩形 60": "较完整流程下降约一个数量级",
    "矩形 65": "100%",
    "矩形 66": "证据损坏切分判定正确率",
    "矩形 67": "均正确判为“证据不足”，最终判定结果不变",
})
# ---- S7 已完成工作总览(4卡) ----
do(6, {
    "矩形 21": "已完成工作总览",
    "矩形: 圆角 366": "① 统一风险图建模",
    "矩形 14": "CDB-RG 八类节点十类边；暴露路径形式化；字段—表—实例三层敏感性聚合（含单调性证明）",
    "矩形: 圆角 368": "② 评测基准 CloudDB-PathBench",
    "矩形 15": "参数化生成框架 + CloudGoat 真实靶场；四种正交切分与多维评测指标",
    "矩形: 圆角 367": "③ EIC-Agent 核心机制",
    "矩形 16": "表达—判定分离；Gate·Score 确定性判定；算子若干性质已形式化",
    "矩形: 圆角 930": "④ 自主探索与提前终止",
    "矩形 17": "代价优先逐维取证；硬维度一票否决提前终止，已完成实测",
})
# ---- S8 EIC-Agent 推理链路(5列同名 koppt-文本框) ----
s8 = p.slides[7]
cols = find_all(s8.shapes, "koppt-文本框", [])
# 按 left 排序保证从左到右
cols.sort(key=lambda sh: sh.left or 0)
col_texts = [
    "候选路径搜索：约束 DFS 在风险图上枚举暴露链",
    "工具化取证：调用 7 类工具沿链收集证据",
    "结构化证据表达：入口/可达/权限/目标/感知 五维向量",
    "Gate·Score 判定：确定性算子给出结论，模型不参与",
    "LLM 归因：生成可回溯的自然语言解释与处置建议",
]
for sh, txt in zip(cols, col_texts):
    fill(sh, txt)
do(7, {
    "矩形 60": "EIC-Agent 推理链路：表达—判定分离",
    "文本框 55": "大模型负责“查证据、说清楚”，确定性算子负责“下判定”。",
})
# ---- S9 统一风险图谱(留给 image②) ----
do(8, {
    "矩形 9": "统一风险图 CDB-RG：系统实测五层暴露图谱（20 节点 / 29 边）",
})
# ---- S10 原型验证与下一步(留给 image①) ----
do(9, {
    "矩形 15": "原型验证与下一步计划",
    "文本框 11": "原型已端到端贯通：场景接入 → 建图 → 取证 → 确定性判定 → 归因 → 可视化",
    "矩形 14": "下一步：① 自主探索完整闭环（信息增益驱动补证/扩展）② 基准规模化与系统性实验 ③ GV-FA 训练与对齐 ④ 论文完善与答辩",
    "矩形 19": "当前定位：方法验证与关键机制打通，真实大规模数据与基线实验推进中",
    "矩形 22": "GV-FA",
    "矩形 24": "下一步",
})
# ---- S11 尾页 ----
do(10, {
    "矩形 21": "谢谢各位专家！　恳请批评指正",
    "文本框 16": "2026 年 7 月    July 2026",
    "文本框 23": "王凌云　导师：蔡忠闽",
})

p.save(OUT)
print(f"✓ 已生成 {OUT}")
