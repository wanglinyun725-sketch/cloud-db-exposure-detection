# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import MSO_AUTO_SIZE

SRC="模板_副本.pptx"; OUT="中期报告_v2.pptx"
IMG1="/Users/yunyun/Desktop/2/20260713175207.jpg"  # 工作台
IMG2="/Users/yunyun/Desktop/2/20260713175154.jpg"  # 风险图谱
p=Presentation(SRC)

def find(shapes,name):
    for sh in shapes:
        if sh.shape_type==6:
            r=find(sh.shapes,name)
            if r is not None: return r
        elif sh.name==name: return sh
    return None
def find_all(shapes,name,acc):
    for sh in shapes:
        if sh.shape_type==6: find_all(sh.shapes,name,acc)
        elif sh.name==name: acc.append(sh)
    return acc

def fill(shape,text,size=None):
    if shape is None: return False
    tf=shape.text_frame
    p0=tf.paragraphs[0]
    if p0.runs:
        p0.runs[0].text=text
        for r in p0.runs[1:]: r._r.getparent().remove(r._r)
    else:
        p0.text=text
    for ex in tf.paragraphs[1:]: ex._p.getparent().remove(ex._p)
    if size and tf.paragraphs[0].runs:
        tf.paragraphs[0].runs[0].font.size=Pt(size)
    try:
        tf.word_wrap=True
        tf.auto_size=MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except: pass
    return True

def do(i,m):
    s=p.slides[i]
    for name,v in m.items():
        if isinstance(v,tuple): txt,sz=v
        else: txt,sz=v,None
        if find(s.shapes,name) is None: print(f"  ⚠ s{i+1} 缺 {name}")
        fill(find(s.shapes,name),txt,sz)

# S1 封面
do(0,{"矩形 21":("面向云数据库高敏数据暴露路径侦测的证据约束智能体方法研究",30),
      "文本框 16":"2026 年 7 月    July 2026",
      "文本框 22":"王凌云　导师：蔡忠闽"})
# S2 目录
do(1,{"矩形 44":("研究背景与目标",22),"矩形 35":("总体方案与路线",22),
      "矩形 28":("已完成工作",22),"矩形 40":("下一步计划",22)})
# S3 背景
do(2,{"矩形 53":("研究背景与问题",54),"矩形 1":("高敏数据核心载体",16),
      "矩形 6":("毒性组合暴露路径",16),"矩形 7":("判定需可审计",16)})
# S4 目标与方案
do(3,{"矩形 22":("研究目标与总体方案",24),
      "矩形 17":("① 评测基准",20),"文本框 16":("CloudDB-PathBench：受控、可复现的暴露路径评测环境",16),
      "矩形 19":("② EIC-Agent",20),"文本框 18":("表达—判定分离，判定可审计、抗幻觉",16),
      "矩形 21":("③ 原型系统",20),"文本框 20":("CloudDB-PathDetect：端到端贯通与可视化验证",16)})
# S5 技术路线
do(4,{"矩形 22":("技术路线总览",36),
      "矩形 24":("统一风险图 CDB-RG",18),"矩形 57":("身份/网络/数据/行为统一为同构风险图",14),
      "矩形 44":("EIC-Agent 智能体",18),"矩形 58":("工具化取证 + 表达—判定分离，判定交确定性算子",14),
      "矩形 51":("提前终止与 GV-FA",18),"矩形 63":("Gate 一票否决剪枝；图验证反馈对齐(路线已定)",14)})
# S6 实测数据
do(5,{"矩形 70":("关键实测：提前终止效率",26),
      "矩形 54":("≈25%",48),"矩形 29":("路径提前终止比例",18),"矩形 31":("68 条候选中 17 条被硬门剪枝",13),
      "矩形 58":("10×",48),"矩形 59":("被终止路径开销下降",18),"矩形 60":("较完整流程降约一个数量级",13),
      "矩形 65":("100%",48),"矩形 66":("证据损坏判定正确率",18),"矩形 67":("均正确判为证据不足，判定不变",13)})
# S7 已完成总览
do(6,{"矩形 21":("已完成工作总览",36),
      "矩形: 圆角 366":("① 风险图建模",16),"矩形 14":("CDB-RG 八类节点十类边；暴露路径形式化；三层敏感性聚合(含单调性证明)",13),
      "矩形: 圆角 368":("② 评测基准",16),"矩形 15":("参数化生成 + CloudGoat 真实靶场；四种正交切分与多维评测指标",13),
      "矩形: 圆角 367":("③ EIC-Agent",16),"矩形 16":("表达—判定分离；Gate·Score 确定性判定；算子性质已形式化",13),
      "矩形: 圆角 930":("④ 提前终止",16),"矩形 17":("代价优先逐维取证；硬维度一票否决提前终止，已完成实测",13)})
# S8 推理链路 → 清空文字列, 放图①
s8=p.slides[7]
for c in find_all(s8.shapes,"koppt-文本框",[]): fill(c,"")
do(7,{"矩形 60":("EIC-Agent 推理链路(工作台实测)",22),"文本框 55":""})
s8.shapes.add_picture(IMG1, Inches(1.66), Inches(1.55), width=Inches(10.0))
# S9 风险图谱 → 放图②
do(8,{"矩形 9":("统一风险图 CDB-RG：实测五层暴露图谱",20)})
p.slides[8].shapes.add_picture(IMG2, Inches(3.66), Inches(1.5), width=Inches(6.0))
# S10 原型与下一步
do(9,{"矩形 15":("原型验证与下一步计划",26),
      "文本框 11":("原型已端到端贯通：场景接入→建图→取证→确定性判定→归因→可视化；当前处于方法验证与机制打通阶段，真实大规模数据与基线实验推进中。",16),
      "矩形 14":("下一步：① 自主探索完整闭环(信息增益驱动补证/扩展)  ② 评测基准规模化与系统性实验  ③ GV-FA 训练与对齐  ④ 论文完善与答辩",16),
      "矩形 19":("进展与规划",20),"矩形 22":("GV-FA",18),"矩形 24":("下一步",18)})
# S11 尾页
do(10,{"矩形 21":("谢谢大家",88),"文本框 16":"2026 年 7 月    July 2026","文本框 23":"王凌云　导师：蔡忠闽"})

p.save(OUT)
print(f"✓ 已生成 {OUT}")
